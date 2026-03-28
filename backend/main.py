from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import os
import shutil
from pathlib import Path
import sys
import requests
import io
import fitz # PyMuPDF
from pptx import Presentation

# Fix relative imports
sys.path.append(os.path.join(os.path.dirname(__file__), "project"))

from project.core.rag_system import RAGSystem
from project.core.document_manager import DocumentManager
import project.config as config

app = FastAPI(title="Drivedrop Agentic RAG")

# Enable CORS (Wildcard Origins require allow_credentials=False)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize RAG System
rag_system = RAGSystem()
rag_system.initialize()
doc_manager = DocumentManager(rag_system)

# --- Models ---

class ChatRequest(BaseModel):
    file_id: Optional[str] = None
    file_name: Optional[str] = "General Chat"
    mime_type: Optional[str] = "text/plain"
    access_token: Optional[str] = None
    query: str
    history: Optional[List[dict]] = []
    thread_id: Optional[str] = None

class ChatResponse(BaseModel):
    answer: str
    thread_id: str

# --- Extraction Helpers ---

def extract_text(content: bytes, mime_type: str) -> str:
    if "pdf" in mime_type:
        doc = fitz.open(stream=content, filetype="pdf")
        return "".join([page.get_text() for page in doc])
    elif "presentation" in mime_type:
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return content.decode("utf-8", errors="ignore")

# --- Endpoints ---

@app.get("/")
def read_root():
    return {"status": "ok", "engine": "Agentic RAG for Dummies"}

@app.post("/chat", response_model=ChatResponse)
async def chat_endpoint(request: ChatRequest):
    try:
        # 1. Handle On-Demand Google Drive Integration
        if request.file_id and request.access_token:
            # Check if we already have this file in our Markdown Knowledge Base
            md_path = Path(config.MARKDOWN_DIR) / f"{request.file_id}.md"
            if not md_path.exists():
                print(f"[Drive-Sync] Fetching new file: {request.file_name}")
                headers = {"Authorization": f"Bearer {request.access_token}"}
                url = f"https://www.googleapis.com/drive/v3/files/{request.file_id}?alt=media"
                response = requests.get(url, headers=headers)
                
                if response.status_code == 200:
                    text = extract_text(response.content, request.mime_type)
                    
                    # Save to Markdown Dir for Persistent Indexing
                    os.makedirs(config.MARKDOWN_DIR, exist_ok=True)
                    with open(md_path, "w", encoding="utf-8") as f:
                        f.write(text)
                    
                    # Tell Doc Manager to index it
                    doc_manager.add_documents([str(md_path)])
                    print(f"[Drive-Sync] Indexed: {request.file_id}")

        # 2. Execute Agentic Reasoning Loop (LangGraph)
        thread_id = request.thread_id or rag_system.thread_id
        run_config = {"configurable": {"thread_id": thread_id}, "recursion_limit": rag_system.recursion_limit}
        
        # Prepare history for LangGraph if needed
        # (Dummies app uses built-in memory, but we can seed it here)
        
        result = rag_system.agent_graph.invoke(
            {"messages": [("user", request.query)]},
            run_config
        )
        
        reply = result["messages"][-1].content
        return ChatResponse(answer=reply, thread_id=thread_id)

    except Exception as e:
        print(f"Error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/clear")
def clear_knowledge():
    doc_manager.clear_all()
    # Also clear temporary markdown cache
    if os.path.exists(config.MARKDOWN_DIR):
        shutil.rmtree(config.MARKDOWN_DIR)
    return {"status": "cleared"}

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 8000)))
